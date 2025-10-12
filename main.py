from flask import Flask, request, send_file, jsonify
import io
import re
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)

# Configure logging
logging.basicConfig(level=logging.DEBUG)

@app.route('/generate-pptx', methods=['POST'])
def generate_pptx():
    try:
        # Parse JSON payload
        data = request.get_json()
        logging.debug(f"Received data: {data}")

        # Extract top-level title for filename
        title = data.get('title', 'lesson_slides')
        # Sanitize filename: replace invalid characters with underscores
        safe_title = re.sub(r'[^a-zA-Zа-яА-Я0-9_]', '_', title.strip())
        filename = f"{safe_title}.pptx"
        logging.debug(f"Generated filename: {filename}")

        # Extract slides data
        slides_data = data.get('slides', [])
        if not slides_data:
            logging.error("No slides data provided")
            return jsonify({'error': 'No slides data provided'}), 400

        # Create PPTX presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[6]  # Blank layout

        for slide_info in slides_data:
            slide = prs.slides.add_slide(title_slide_layout)

            # Slide title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = slide_info.get('title', 'Untitled Slide')
            title_p = title_frame.paragraphs[0]
            title_p.font.size = Pt(32)
            title_p.alignment = PP_ALIGN.CENTER

            # Bullet points
            left = Inches(0.8)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            body_shape = slide.shapes.add_textbox(left, top, width, height)
            body_frame = body_shape.text_frame
            for bullet in slide_info.get('bullets', []):
                p = body_frame.add_paragraph()
                p.text = '• ' + bullet
                p.font.size = Pt(18)
                p.level = 0

            # Notes
            if slide_info.get('note'):
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.add_paragraph().text = slide_info['note']

        # Save to BytesIO
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Return PPTX file with dynamic filename
        return send_file(
            pptx_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        logging.error(f"Error generating PPTX: {str(e)}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)